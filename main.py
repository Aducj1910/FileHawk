"""
FileFinder - Local Semantic File Search Tool

REQUIRED DEPENDENCIES:
- sentence-transformers
- chromadb
- pdfplumber

OPTIONAL DEPENDENCIES (for extended file type support):
- python-docx (for .docx files)
- pandas (for .csv and .xlsx files)
- openpyxl (for .xlsx files - alternative to pandas)
- beautifulsoup4 (for .html files)

Install with: pip install sentence-transformers chromadb pdfplumber python-docx pandas openpyxl beautifulsoup4
"""

import os
import csv
import json
import time
import sys
import re
import hashlib
from collections import Counter, defaultdict
from math import log, sqrt
import numpy as np
import pdfplumber
from sentence_transformers import SentenceTransformer
import chromadb
from metadata_tracker import MetadataTracker
from folder_settings import folder_settings

# Import configuration settings
try:
    from config import *
except ImportError:
    print("ERROR: config.py file not found!")
    print("Please ensure config.py exists in the same directory as main.py")
    sys.exit(1)

# Terminal colors and formatting
class Colors:
    HEADER = '\033[95m'
    OKBLUE = '\033[94m'
    OKCYAN = '\033[96m'
    OKGREEN = '\033[92m'
    WARNING = '\033[93m'
    FAIL = '\033[91m'
    ENDC = '\033[0m'
    BOLD = '\033[1m'
    UNDERLINE = '\033[4m'

def print_header(text):
    """Print a formatted header"""
    try:
        print(f"\n{Colors.HEADER}{Colors.BOLD}{'='*60}")
        print(f"  {text}")
        print(f"{'='*60}{Colors.ENDC}\n")
    except UnicodeEncodeError:
        print(f"\n{'='*60}")
        print(f"  {text}")
        print(f"{'='*60}\n")

def print_step(text):
    """Print a step with blue color"""
    try:
        print(f"{Colors.OKBLUE}> {text}{Colors.ENDC}")
    except UnicodeEncodeError:
        # Fallback for Windows console encoding issues
        print(f"> {text}")

def print_success(text):
    """Print success message in green"""
    try:
        print(f"{Colors.OKGREEN}+ {text}{Colors.ENDC}")
    except UnicodeEncodeError:
        print(f"+ {text}")

def print_warning(text):
    """Print warning in yellow"""
    try:
        print(f"{Colors.WARNING}! {text}{Colors.ENDC}")
    except UnicodeEncodeError:
        print(f"! {text}")

def print_error(text):
    """Print error in red"""
    try:
        print(f"{Colors.FAIL}- {text}{Colors.ENDC}")
    except UnicodeEncodeError:
        print(f"- {text}")

def loading_animation(text, duration=None):
    """Show a loading animation"""
    if not SHOW_LOADING_ANIMATIONS:
        return
    
    if duration is None:
        duration = LOADING_ANIMATION_DURATION
    
    # Cross-platform animation characters
    animation = ["|", "/", "-", "\\"]
    start_time = time.time()
    i = 0
    
    try:
        while time.time() - start_time < duration:
            sys.stdout.write(f"\r{Colors.OKCYAN}{animation[i % len(animation)]} {text}{Colors.ENDC}")
            sys.stdout.flush()
            time.sleep(0.25)
            i += 1
        
        sys.stdout.write(f"\r{' ' * (len(text) + 2)}\r")  # Clear the line
        sys.stdout.flush()
    except UnicodeEncodeError:
        # Fallback for Windows console encoding issues
        while time.time() - start_time < duration:
            sys.stdout.write(f"\r{animation[i % len(animation)]} {text}")
            sys.stdout.flush()
            time.sleep(0.25)
            i += 1
        
        sys.stdout.write(f"\r{' ' * (len(text) + 2)}\r")  # Clear the line
        sys.stdout.flush()

def progress_bar(current, total, prefix="Progress"):
    """Show a progress bar"""
    if not SHOW_PROGRESS_BARS:
        return
    
    bar_length = 40
    filled_length = int(round(bar_length * current / float(total)))
    percents = round(100.0 * current / float(total), 1)
    bar = '#' * filled_length + '-' * (bar_length - filled_length)
    sys.stdout.write(f'\r{Colors.OKCYAN}{prefix}: [{bar}] {percents}% ({current}/{total}){Colors.ENDC}')
    sys.stdout.flush()
    if current == total:
        print()  # New line when complete

# Additional imports for file type handling
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print_warning("python-docx not installed. .docx files will be skipped.")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    print_warning("pandas not installed. .csv and .xlsx files will be skipped.")

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    print_warning("openpyxl not installed. .xlsx files will be skipped.")

try:
    from bs4 import BeautifulSoup
    BEAUTIFULSOUP_AVAILABLE = True
except ImportError:
    BEAUTIFULSOUP_AVAILABLE = False
    print_warning("beautifulsoup4 not installed. .html files will be skipped.")

# RTF support
try:
    import striprtf
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False
    print_warning("striprtf not installed. .rtf files will be skipped.")

# LibreOffice/OpenOffice support
try:
    from odf import text, teletype
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False
    print_warning("odfpy not installed. .odt, .ods, .odp files will be skipped.")

# PowerPoint support
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print_warning("python-pptx not installed. .pptx files will be skipped.")

# E-book support
try:
    import ebooklib
    from ebooklib import epub
    EPUB_AVAILABLE = True
except ImportError:
    EPUB_AVAILABLE = False
    print_warning("ebooklib not installed. .epub files will be skipped.")

# Email support
try:
    import email
    from email import policy
    EMAIL_AVAILABLE = True
except ImportError:
    EMAIL_AVAILABLE = False
    print_warning("email module not available. .eml files will be skipped.")

try:
    import extract_msg
    MSG_AVAILABLE = True
except ImportError:
    MSG_AVAILABLE = False
    print_warning("extract-msg not installed. .msg files will be skipped.")



# Supported file extensions with special handling
SPECIAL_EXTENSIONS = {
    '.pdf': 'pdfplumber',
    '.docx': 'python-docx', 
    '.csv': 'csv module',
    '.xlsx': 'openpyxl/pandas',
    '.html': 'BeautifulSoup',
    '.json': 'json module',
    '.rtf': 'striprtf',
    '.odt': 'odfpy',
    '.ods': 'odfpy',
    '.odp': 'odfpy',
    '.pptx': 'python-pptx',
    '.xml': 'xml parser',
    '.yaml': 'yaml parser',
    '.yml': 'yaml parser',
    '.ini': 'ini parser',
    '.epub': 'ebooklib',
    '.eml': 'email parser',
    '.msg': 'extract-msg'
}

def show_supported_formats():
    """Display supported file formats and their parsers"""
    print_header("FileFinder - Supported File Formats")
    
    print_step("Special handling (with specific parsers):")
    for ext, parser in SPECIAL_EXTENSIONS.items():
        # Determine status for each file type
        status = "?"
        if ext == '.pdf':
            status = "+"
        elif ext == '.json':
            status = "+"
        elif ext == '.csv':
            status = "+"
        elif ext == '.docx':
            status = "+" if DOCX_AVAILABLE else "-"
        elif ext == '.xlsx':
            status = "+" if (PANDAS_AVAILABLE or OPENPYXL_AVAILABLE) else "-"
        elif ext == '.html':
            status = "+" if BEAUTIFULSOUP_AVAILABLE else "-"
        elif ext == '.rtf':
            status = "+" if RTF_AVAILABLE else "-"
        elif ext in ['.odt', '.ods', '.odp']:
            status = "+" if ODF_AVAILABLE else "-"
        elif ext == '.pptx':
            status = "+" if PPTX_AVAILABLE else "-"
        elif ext in ['.xml', '.yaml', '.yml', '.ini']:
            status = "+"  # Built-in Python modules
        elif ext == '.epub':
            status = "+" if EPUB_AVAILABLE else "-"
        elif ext == '.eml':
            status = "+" if EMAIL_AVAILABLE else "-"
        elif ext == '.msg':
            status = "+" if MSG_AVAILABLE else "-"
        
        status_color = Colors.OKGREEN if status == "+" else Colors.FAIL
        print(f"  {Colors.BOLD}{ext:8}{Colors.ENDC} - {parser:15} [{status_color}{status}{Colors.ENDC}]")
    
    print(f"\n{Colors.BOLD}All other file types (code, config, data files):{Colors.ENDC}")
    print(f"  {Colors.OKCYAN}.py, .js, .ts, .tsx, .c, .cpp, .h, .sh, .ini, .yaml, etc.{Colors.ENDC}")
    print(f"  {Colors.WARNING}-> Treated as plain text (UTF-8 with error handling){Colors.ENDC}")
    print(f"\n{Colors.WARNING}Note: Files starting with '.' or '~' are automatically skipped.{Colors.ENDC}")
    print()

# ---- INIT ----
print_header("FileFinder - Initialization")

# Global model loading status
model_loading_status = {
    "is_loading": True,
    "progress": 0,
    "message": "Initializing...",
    "error": None
}

print_step("Loading embedding model...")
model_loading_status.update({
    "message": "Loading sentence transformer model...",
    "progress": 10
})
loading_animation("Loading sentence transformer model", 1.5)

try:
    model_loading_status.update({
        "message": "Downloading model weights...",
        "progress": 30
    })
    model = SentenceTransformer(EMBEDDING_MODEL)
    model_loading_status.update({
        "is_loading": False,
        "progress": 100,
        "message": f"Model loaded successfully: {EMBEDDING_MODEL}",
        "error": None
    })
    print_success(f"Embedding model loaded: {EMBEDDING_MODEL}")
except Exception as e:
    model_loading_status.update({
        "is_loading": False,
        "progress": 0,
        "message": f"Failed to load model: {str(e)}",
        "error": str(e)
    })
    print_error(f"Failed to load embedding model: {e}")
    raise

print_step("starting metadata tracker...")
model_loading_status.update({
    "message": "Setting up metadata tracking...",
    "progress": 55
})
loading_animation("Setting up metadata tracking", 1)
metadata_tracker = MetadataTracker(APP_NAME)
print_success("metadata tracker ready")

print_step("starting chromadb client...")
model_loading_status.update({
    "message": "Connecting to vector database...",
    "progress": 60
})
loading_animation("Connecting to vector database", 1)
db_dir = metadata_tracker.get_db_directory()
# Use PersistentClient for ChromaDB 1.0+
client = chromadb.PersistentClient(path=str(db_dir))
print_success(f"chromadb client ready with storage at: {db_dir}")

# Collections setup based on chunking configuration
if ENABLE_GRANULAR_CHUNKING:
    print_step("Setting up dual-chunk collections...")
    model_loading_status.update({
        "message": "Setting up dual-chunk collections...",
        "progress": 80
    })
    granular_collection = client.get_or_create_collection("filefinder_granular")  # Line-by-line
    filelevel_collection = client.get_or_create_collection("filefinder_filelevel")  # Large chunks
    print_success("Dual-chunk collections ready")
else:
    print_step("Setting up file-level collection...")
    model_loading_status.update({
        "message": "Setting up file-level collection...",
        "progress": 80
    })
    granular_collection = None  # Not used in file-level only mode
    filelevel_collection = client.get_or_create_collection("filefinder_filelevel")  # Large chunks only
    print_success("File-level collection ready")

model_loading_status.update({
    "message": "Initialization complete",
    "progress": 100
})

def extract_text(file_path):
    """
    Extract text from various file types using appropriate parsers.
    Returns empty string if extraction fails.
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        # Special file type handling
        if file_ext == '.pdf':
            return extract_pdf_text(file_path)
        elif file_ext == '.docx':
            return extract_docx_text(file_path)
        elif file_ext == '.csv':
            return extract_csv_text(file_path)
        elif file_ext == '.xlsx':
            return extract_xlsx_text(file_path)
        elif file_ext == '.html':
            return extract_html_text(file_path)
        elif file_ext == '.json':
            return extract_json_text(file_path)
        elif file_ext == '.rtf':
            return extract_rtf_text(file_path)
        elif file_ext in ['.odt', '.ods', '.odp']:
            return extract_odf_text(file_path)
        elif file_ext == '.pptx':
            return extract_pptx_text(file_path)
        elif file_ext in ['.xml', '.yaml', '.yml', '.ini']:
            return extract_config_text(file_path, file_ext)
        elif file_ext == '.epub':
            return extract_epub_text(file_path)
        elif file_ext == '.eml':
            return extract_eml_text(file_path)
        elif file_ext == '.msg':
            return extract_msg_text(file_path)
        else:
            # All other file types - treat as plain text
            return extract_plain_text(file_path)
            
    except Exception as e:
        print_error(f"Error reading file {file_path}: {e}")
        return ""

def extract_pdf_text(file_path):
    """Extract text from PDF using pdfplumber"""
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            text += page_text + "\n"
    print_success(f"Extracted {len(text)} chars from PDF: {os.path.basename(file_path)}")
    return text

def extract_docx_text(file_path):
    """Extract text from DOCX using python-docx"""
    if not DOCX_AVAILABLE:
        print_error(f"python-docx not available for {os.path.basename(file_path)}")
        return ""
    
    try:
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
            text += "\n"
        
        print_success(f"Extracted {len(text)} chars from DOCX: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error extracting DOCX text from {os.path.basename(file_path)}: {e}")
        # Return a placeholder instead of empty string to prevent deletion
        return f"[Error reading DOCX file: {os.path.basename(file_path)}]"

def extract_csv_text(file_path):
    """Extract text from CSV using pandas or csv module"""
    if PANDAS_AVAILABLE:
        # Use pandas for better handling of various CSV formats
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
    else:
        # Fallback to csv module
        text = ""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                text += " ".join(row) + "\n"
    print_success(f"Extracted {len(text)} chars from CSV: {os.path.basename(file_path)}")
    return text

def extract_xlsx_text(file_path):
    """Extract text from XLSX using openpyxl or pandas"""
    if PANDAS_AVAILABLE:
        # Use pandas for better handling - read all sheets
        excel_file = pd.ExcelFile(file_path)
        all_text = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            # Drop rows where all values are NaN (empty rows)
            df = df.dropna(how='all')
            # Also drop columns where all values are NaN (empty columns)
            df = df.dropna(axis=1, how='all')
            
            # Only add sheet if it has content
            if not df.empty:
                sheet_text = df.to_string(index=False)
                if sheet_text.strip():
                    all_text.append(f"Sheet: {sheet_name}\n{sheet_text}")
        
        text = "\n\n".join(all_text)
    elif OPENPYXL_AVAILABLE:
        # Fallback to openpyxl
        wb = load_workbook(file_path, data_only=True)
        all_text = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_rows = []
            
            for row in ws.iter_rows(values_only=True):
                # Skip rows where all cells are None (empty rows)
                if not all(cell is None for cell in row):
                    # Only include non-None cells in the text
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():  # Only add non-empty row text
                        sheet_rows.append(row_text)
            
            # Only add sheet if it has content
            if sheet_rows:
                sheet_text = "\n".join(sheet_rows)
                all_text.append(f"Sheet: {sheet_name}\n{sheet_text}")
        
        text = "\n\n".join(all_text)
    else:
        print_error(f"Neither pandas nor openpyxl available for {os.path.basename(file_path)}")
        return ""
    print_success(f"Extracted {len(text)} chars from XLSX: {os.path.basename(file_path)}")
    return text

def extract_html_text(file_path):
    """Extract visible text from HTML using BeautifulSoup"""
    if not BEAUTIFULSOUP_AVAILABLE:
        print_error(f"beautifulsoup4 not available for {os.path.basename(file_path)}")
        return ""
    
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    
    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.decompose()
    
    # Get text and clean up whitespace
    text = soup.get_text()
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    text = ' '.join(chunk for chunk in chunks if chunk)
    
    print_success(f"Extracted {len(text)} chars from HTML: {os.path.basename(file_path)}")
    return text

def extract_json_text(file_path):
    """Extract and flatten text from JSON"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        data = json.load(f)
    
    def flatten_json(obj, text_parts=None):
        if text_parts is None:
            text_parts = []
        
        if isinstance(obj, dict):
            for key, value in obj.items():
                text_parts.append(str(key))
                flatten_json(value, text_parts)
        elif isinstance(obj, list):
            for item in obj:
                flatten_json(item, text_parts)
        else:
            text_parts.append(str(obj))
        
        return text_parts
    
    text_parts = flatten_json(data)
    text = " ".join(text_parts)
    print_success(f"Extracted {len(text)} chars from JSON: {os.path.basename(file_path)}")
    return text

def extract_rtf_text(file_path):
    """Extract text from RTF using striprtf"""
    if not RTF_AVAILABLE:
        print_error(f"striprtf not available for {os.path.basename(file_path)}")
        return ""
    
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        rtf_content = f.read()
    
    text = rtf_to_text(rtf_content)
    print_success(f"Extracted {len(text)} chars from RTF: {os.path.basename(file_path)}")
    return text

def extract_odf_text(file_path):
    """Extract text from LibreOffice/OpenOffice documents using odfpy"""
    if not ODF_AVAILABLE:
        print_error(f"odfpy not available for {os.path.basename(file_path)}")
        return ""
    
    try:
        text = teletype.text(file_path)
        print_success(f"Extracted {len(text)} chars from ODF: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error extracting ODF text: {e}")
        return ""

def extract_pptx_text(file_path):
    """Extract text from PowerPoint using python-pptx"""
    if not PPTX_AVAILABLE:
        print_error(f"python-pptx not available for {os.path.basename(file_path)}")
        return ""
    
    try:
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        print_success(f"Extracted {len(text)} chars from PowerPoint: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error extracting PowerPoint text: {e}")
        return ""

def extract_config_text(file_path, file_ext):
    """Extract text from config files (XML, YAML, INI)"""
    try:
        if file_ext == '.xml':
            return extract_xml_text(file_path)
        elif file_ext in ['.yaml', '.yml']:
            return extract_yaml_text(file_path)
        elif file_ext == '.ini':
            return extract_ini_text(file_path)
        else:
            return extract_plain_text(file_path)
    except Exception as e:
        print_error(f"Error extracting config text: {e}")
        return extract_plain_text(file_path)

def extract_xml_text(file_path):
    """Extract and flatten text from XML"""
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        def extract_xml_elements(element, text_parts=None):
            if text_parts is None:
                text_parts = []
            
            # Add tag name
            text_parts.append(element.tag)
            
            # Add attributes
            for key, value in element.attrib.items():
                text_parts.append(f"{key}={value}")
            
            # Add text content
            if element.text and element.text.strip():
                text_parts.append(element.text.strip())
            
            # Recursively process children
            for child in element:
                extract_xml_elements(child, text_parts)
            
            return text_parts
        
        text_parts = extract_xml_elements(root)
        text = " ".join(text_parts)
        print_success(f"Extracted {len(text)} chars from XML: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error parsing XML: {e}")
        return extract_plain_text(file_path)

def extract_yaml_text(file_path):
    """Extract and flatten text from YAML"""
    try:
        import yaml
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = yaml.safe_load(f)
        
        def flatten_yaml(obj, text_parts=None):
            if text_parts is None:
                text_parts = []
            
            if isinstance(obj, dict):
                for key, value in obj.items():
                    text_parts.append(str(key))
                    flatten_yaml(value, text_parts)
            elif isinstance(obj, list):
                for item in obj:
                    flatten_yaml(item, text_parts)
            else:
                text_parts.append(str(obj))
            
            return text_parts
        
        text_parts = flatten_yaml(data)
        text = " ".join(text_parts)
        print_success(f"Extracted {len(text)} chars from YAML: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error parsing YAML: {e}")
        return extract_plain_text(file_path)

def extract_ini_text(file_path):
    """Extract text from INI files"""
    try:
        import configparser
        config = configparser.ConfigParser()
        config.read(file_path, encoding='utf-8')
        
        text_parts = []
        for section in config.sections():
            text_parts.append(f"[{section}]")
            for key, value in config[section].items():
                text_parts.append(f"{key}={value}")
        
        text = " ".join(text_parts)
        print_success(f"Extracted {len(text)} chars from INI: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error parsing INI: {e}")
        return extract_plain_text(file_path)

def extract_epub_text(file_path):
    """Extract text from EPUB using ebooklib"""
    if not EPUB_AVAILABLE:
        print_error(f"ebooklib not available for {os.path.basename(file_path)}")
        return ""
    
    try:
        book = epub.read_epub(file_path)
        text = ""
        
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                # Extract text from HTML content
                content = item.get_content().decode('utf-8')
                # Simple HTML stripping
                import re
                text += re.sub('<[^<]+?>', '', content) + "\n"
        
        print_success(f"Extracted {len(text)} chars from EPUB: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error extracting EPUB text: {e}")
        return ""

def extract_eml_text(file_path):
    """Extract text from EML files using email module"""
    if not EMAIL_AVAILABLE:
        print_error(f"email module not available for {os.path.basename(file_path)}")
        return ""
    
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            msg = email.message_from_file(f, policy=policy.default)
        
        text = ""
        
        # Extract headers
        for header, value in msg.items():
            text += f"{header}: {value}\n"
        
        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    text += part.get_content() + "\n"
        else:
            text += msg.get_content() + "\n"
        
        print_success(f"Extracted {len(text)} chars from EML: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error extracting EML text: {e}")
        return ""

def extract_msg_text(file_path):
    """Extract text from MSG files using extract-msg"""
    if not MSG_AVAILABLE:
        print_error(f"extract-msg not available for {os.path.basename(file_path)}")
        return ""
    
    try:
        msg = extract_msg.Message(file_path)
        
        text = ""
        
        # Extract headers
        text += f"From: {msg.sender}\n"
        text += f"To: {msg.to}\n"
        text += f"Subject: {msg.subject}\n"
        text += f"Date: {msg.date}\n\n"
        
        # Extract body
        text += msg.body + "\n"
        
        print_success(f"Extracted {len(text)} chars from MSG: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print_error(f"Error extracting MSG text: {e}")
        return ""

def extract_plain_text(file_path):
    """Extract text from any file type as plain text"""
    with open(file_path, "r", encoding=DEFAULT_ENCODING, errors=ENCODING_ERROR_HANDLING) as f:
        text = f.read()
    print_success(f"Extracted {len(text)} chars from plain text: {os.path.basename(file_path)}")
    return text

# ============================================================================
# ENHANCED GIST MODE UTILITIES
# ============================================================================

def clean_text_for_gist(text):
    """Clean text by removing boilerplate and normalizing whitespace"""
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # Strip whitespace
        cleaned = line.strip()
        
        # Skip empty lines and common boilerplate patterns
        if not cleaned:
            continue
        if re.match(r'^[-=_*#]{3,}$', cleaned):  # Separator lines
            continue
        if re.match(r'^\s*(copyright|©|\(c\))', cleaned, re.IGNORECASE):
            continue
        if re.match(r'^\s*generated\s+(automatically|by)', cleaned, re.IGNORECASE):
            continue
        
        # Collapse repeated whitespace
        cleaned = re.sub(r'\s+', ' ', cleaned)
        
        if cleaned:
            cleaned_lines.append(cleaned)
    
    return cleaned_lines

def cosine_similarity(vec1, vec2):
    """Compute cosine similarity between two vectors"""
    if len(vec1) != len(vec2):
        return 0.0
    
    dot_product = sum(a * b for a, b in zip(vec1, vec2))
    norm_a = sqrt(sum(a * a for a in vec1))
    norm_b = sqrt(sum(b * b for b in vec2))
    
    if norm_a == 0 or norm_b == 0:
        return 0.0
    
    return dot_product / (norm_a * norm_b)

def compute_tf_idf_terms(text_chunks, max_terms=50):
    """Compute top TF-IDF terms for a set of text chunks"""
    # Simple tokenization
    def tokenize(text):
        # Extract unigrams and bigrams
        words = re.findall(r'\b\w{2,}\b', text.lower())
        unigrams = words
        bigrams = [f"{words[i]}_{words[i+1]}" for i in range(len(words)-1)]
        return unigrams + bigrams
    
    # Compute term frequencies per chunk
    chunk_terms = []
    all_terms = set()
    
    for chunk in text_chunks:
        terms = tokenize(chunk)
        chunk_term_counts = Counter(terms)
        chunk_terms.append(chunk_term_counts)
        all_terms.update(terms)
    
    if not all_terms:
        return []
    
    # Compute document frequency for each term
    df = {}
    for term in all_terms:
        df[term] = sum(1 for chunk_counts in chunk_terms if term in chunk_counts)
    
    # Compute TF-IDF scores
    num_chunks = len(text_chunks)
    term_scores = defaultdict(float)
    
    for chunk_counts in chunk_terms:
        chunk_len = sum(chunk_counts.values()) or 1
        
        for term, count in chunk_counts.items():
            tf = count / chunk_len
            idf = log(num_chunks / df[term]) if df[term] > 0 else 0
            term_scores[term] += tf * idf
    
    # Return top terms sorted by score
    sorted_terms = sorted(term_scores.items(), key=lambda x: x[1], reverse=True)
    return [term for term, score in sorted_terms[:max_terms]]

def normalize_embedding(embedding):
    """Normalize an embedding vector to unit length"""
    if isinstance(embedding, list):
        embedding = np.array(embedding)
    
    norm = np.linalg.norm(embedding)
    if norm == 0:
        return embedding
    
    return embedding / norm

def compute_file_metadata_gist(chunks, chunk_embeddings, file_path):
    """
    Compute per-file metadata aggregates for gist mode.
    Returns dict with centroid, maxpool, top_terms, positions, etc.
    """
    if not chunks or not chunk_embeddings:
        return None
    
    try:
        # Convert embeddings to numpy arrays if needed
        embeddings_array = np.array(chunk_embeddings)
        
        # Compute file centroid (mean of normalized embeddings, then re-normalize)
        file_centroid = np.mean(embeddings_array, axis=0)
        file_centroid = normalize_embedding(file_centroid)
        
        # Compute file maxpool (element-wise max, then re-normalize)
        file_maxpool = np.max(embeddings_array, axis=0)
        file_maxpool = normalize_embedding(file_maxpool)
        
        # Extract chunk texts for TF-IDF
        chunk_texts = [chunk_text for chunk_text, _ in chunks]
        
        # Compute top terms using TF-IDF
        top_terms = compute_tf_idf_terms(chunk_texts, GIST_MAX_TOP_TERMS)
        
        # Compute positions (start line numbers for each chunk)
        positions = []
        for _, line_ranges in chunks:
            if line_ranges:
                positions.append(min(line_ranges))
            else:
                positions.append(0)
        
        # Get file metadata
        try:
            stat = os.stat(file_path)
            modified_time = stat.st_mtime
            file_size_bytes = stat.st_size
        except OSError:
            modified_time = time.time()
            file_size_bytes = 0
        
        metadata = {
            'n_chunks': len(chunks),
            'file_centroid': file_centroid.tolist(),
            'file_maxpool': file_maxpool.tolist(),
            'top_terms': top_terms,
            'positions': positions,
            'modified_time': modified_time,
            'file_size_bytes': file_size_bytes,
            'gist_version': '1.0'  # For future migrations
        }
        
        return metadata
        
    except Exception as e:
        print_error(f"Error computing file metadata for {file_path}: {e}")
        return None

def create_large_chunks(content, file_ext):
    """
    Create larger chunks for file-level search.
    Uses paragraph-based chunking with size limits.
    """
    lines = [line.strip() for line in content.split('\n') if line.strip()]
    if not lines:
        return []
    
    chunks = []
    current_chunk = []
    current_line_nums = []
    
    for i, line in enumerate(lines):
        current_chunk.append(line)
        current_line_nums.append(i+1)
        
        # Start new chunk on paragraph breaks (empty lines) or size limits
        if (not line.strip() and current_chunk) or len(current_chunk) >= FILE_LEVEL_CHUNK_SIZE:
            if current_chunk:
                chunks.append(('\n'.join(current_chunk), current_line_nums))
            current_chunk = []
            current_line_nums = []
    
    # Add remaining chunk
    if current_chunk:
        chunks.append(('\n'.join(current_chunk), current_line_nums))
    
    return chunks

def create_gist_chunks(content, file_ext):
    """
    Create enhanced gist chunks for topic-level understanding.
    Features:
    - Text cleaning and boilerplate removal
    - Overlapping chunks to reduce boundary loss
    - Near-duplicate suppression
    - Paragraph/section level chunking
    """
    # Clean the text first
    cleaned_lines = clean_text_for_gist(content)
    if not cleaned_lines:
        return []
    
    chunks = []
    chunk_embeddings = []  # For deduplication
    
    i = 0
    while i < len(cleaned_lines):
        current_chunk = []
        current_line_nums = []
        
        # Build chunk up to size limit
        for j in range(i, min(i + GIST_CHUNK_SIZE, len(cleaned_lines))):
            current_chunk.append(cleaned_lines[j])
            current_line_nums.append(j + 1)
        
        if not current_chunk:
            break
            
        chunk_text = '\n'.join(current_chunk)
        
        # Skip very short chunks
        if len(chunk_text.strip()) < MIN_CONTENT_LENGTH:
            i += max(1, GIST_CHUNK_SIZE - GIST_CHUNK_OVERLAP)
            continue
        
        # Near-duplicate suppression if enabled
        should_add = True
        if GIST_ENABLE_DEDUPLICATION and chunk_embeddings:
            try:
                # Use a simple model for deduplication (this could be optimized)
                from sentence_transformers import SentenceTransformer
                if 'dedup_model' not in globals():
                    global dedup_model
                    dedup_model = SentenceTransformer('all-MiniLM-L6-v2')  # Lightweight model for deduplication
                
                current_embedding = normalize_embedding(dedup_model.encode(chunk_text))
                
                # Check similarity with previous chunks
                for prev_embedding in chunk_embeddings[-5:]:  # Only check last 5 chunks for efficiency
                    similarity = cosine_similarity(current_embedding, prev_embedding)
                    if similarity > GIST_DEDUP_THRESHOLD:
                        should_add = False
                        break
                
                if should_add:
                    chunk_embeddings.append(current_embedding)
                    
            except Exception as e:
                print_warning(f"Deduplication failed, including chunk anyway: {e}")
                should_add = True
        
        if should_add:
            chunks.append((chunk_text, current_line_nums))
        
        # Move to next chunk with overlap
        i += max(1, GIST_CHUNK_SIZE - GIST_CHUNK_OVERLAP)
    
    return chunks

def create_pinpoint_chunks(content, file_ext):
    """
    Create pinpoint chunks for exact phrase matching.
    Uses smaller chunks (sentence/line level) for precision.
    """
    lines = [line.strip() for line in content.split('\n') if line.strip()]
    if not lines:
        return []
    
    chunks = []
    current_chunk = []
    current_line_nums = []
    
    for i, line in enumerate(lines):
        current_chunk.append(line)
        current_line_nums.append(i+1)
        
        # Start new chunk on sentence boundaries or size limits
        sentence_end = line.strip().endswith(('.', '!', '?', ';', ':'))
        if sentence_end or len(current_chunk) >= PINPOINT_CHUNK_SIZE:
            if current_chunk:
                chunks.append(('\n'.join(current_chunk), current_line_nums))
            current_chunk = []
            current_line_nums = []
    
    # Add remaining chunk
    if current_chunk:
        chunks.append(('\n'.join(current_chunk), current_line_nums))
    
    return chunks

def index_folders(folders, chunking_mode='gist'):
    """Index multiple folders with metadata tracking to skip unchanged files"""
    if not folders:
        print_error("no folders to index")
        return
    
    print_header(f"filefinder - {chunking_mode} mode indexing")
    print_step(f"indexing {len(folders)} folder(s)")
    
    for i, folder in enumerate(folders, 1):
        print_step(f"folder {i}/{len(folders)}: {os.path.abspath(folder)}")
        if not os.path.exists(folder):
            print_error(f"folder does not exist: {folder}")
            continue
        
        index_single_folder(folder, chunking_mode)
        print()

def index_single_folder(folder, chunking_mode='gist'):
    """Index a single folder with metadata tracking to skip unchanged files"""
    if ENABLE_GRANULAR_CHUNKING:
        print_header("FileFinder - Dual-Chunk Indexing")
    else:
        print_header("FileFinder - File-Level Indexing")
    print_step(f"Indexing files in: {os.path.abspath(folder)}")
    
    # First, scan for files
    all_files = []
    for root, dirs, files in os.walk(folder):
        for fname in files:
            all_files.append((root, fname))
    
    print_success(f"Found {len(all_files)} files to process")
    
    # Get file paths for metadata checking
    file_paths = [os.path.join(root, fname) for root, fname in all_files]
    
    # Check which files need indexing vs can be skipped
    files_to_index, files_to_skip = metadata_tracker.get_files_to_index(file_paths, chunking_mode)
    
    if files_to_skip:
        print_success(f"{len(files_to_skip)} files unchanged, skipping re-indexing")
    
    if files_to_index:
        print_step(f"Processing {len(files_to_index)} files that need indexing...")
    else:
        print_success("All files are up to date!")
        return
    
    print()
    
    files_indexed = 0
    granular_chunks = 0
    filelevel_chunks = 0
    skipped_files = len(files_to_skip)
    
    # Process only files that need indexing
    for file_idx, (root, fname) in enumerate(all_files):
        fpath = os.path.join(root, fname)
        
        # Skip if file doesn't need indexing
        if fpath not in files_to_index:
            continue
            
        file_ext = os.path.splitext(fname)[1].lower()
        
        # Skip hidden/system files if configured
        if SKIP_HIDDEN_FILES and fname.startswith('.'):
            print_warning(f"Skipped (hidden): {fname}")
            skipped_files += 1
            continue
        if SKIP_SYSTEM_FILES and fname.startswith('~'):
            print_warning(f"Skipped (system): {fname}")
            skipped_files += 1
            continue
        
        # Check file size if limit is set
        if MAX_FILE_SIZE_MB > 0:
            try:
                file_size_mb = os.path.getsize(fpath) / (1024 * 1024)
                if file_size_mb > MAX_FILE_SIZE_MB:
                    print_warning(f"Skipped (too large: {file_size_mb:.1f}MB): {fname}")
                    skipped_files += 1
                    continue
            except OSError:
                pass
        
        print(f"{Colors.BOLD}[{file_idx+1}/{len(all_files)}]{Colors.ENDC} Processing: {Colors.OKCYAN}{fname}{Colors.ENDC}")
        
        content = extract_text(fpath)
        if not content.strip() or len(content.strip()) < MIN_CONTENT_LENGTH:
            print_warning(f"Skipped (empty/too short): {fname}")
            skipped_files += 1
            continue
        
        # Split content into non-empty lines for granular chunking
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        if not lines:
            print_warning(f"No non-empty lines in: {fname}")
            skipped_files += 1
            continue
        
        # Index granular chunks (line-by-line) - only if enabled
        if ENABLE_GRANULAR_CHUNKING:
            print_step(f"  Creating granular chunks (line-by-line)...")
            progress_bar(0, len(lines), "Granular indexing")
            for i, line in enumerate(lines):
                embedding = model.encode(line)
                unique_id = f"granular-{fpath}-{i+1}"
                granular_collection.add(
                    embeddings=[embedding],
                    documents=[line],
                    metadatas=[{"path": fpath, "fname": fname, "line_num": i+1, "file_type": file_ext}],
                    ids=[unique_id]
                )
                granular_chunks += 1
                progress_bar(i+1, len(lines), "Granular indexing")
        
        # Index file-level chunks (large chunks)
        print_step(f"  Creating file-level chunks (large sections)...")
        large_chunks = create_large_chunks(content, file_ext)
        progress_bar(0, len(large_chunks), "File-level indexing")
        for chunk_idx, (chunk_text, line_ranges) in enumerate(large_chunks):
            embedding = model.encode(chunk_text)
            unique_id = f"filelevel-{fpath}-{chunk_idx+1}"
            filelevel_collection.add(
                embeddings=[embedding],
                documents=[chunk_text],
                metadatas=[{
                    "path": fpath, 
                    "fname": fname, 
                    "file_type": file_ext,
                    "chunk_id": chunk_idx + 1,
                    "line_ranges": str(line_ranges),
                    "chunk_size": len(line_ranges)
                }],
                ids=[unique_id]
            )
            filelevel_chunks += 1
            progress_bar(chunk_idx+1, len(large_chunks), "File-level indexing")
        
        if ENABLE_GRANULAR_CHUNKING:
            print_success(f"{fname} | Granular: {len(lines)} chunks | File-level: {len(large_chunks)} chunks | Type: {file_ext}")
        else:
            print_success(f"{fname} | File-level: {len(large_chunks)} chunks | Type: {file_ext}")
        
        # Update metadata for this file
        total_chunks = (len(lines) if ENABLE_GRANULAR_CHUNKING else 0) + len(large_chunks)
        chunk_sizes = [len(line_ranges) for _, line_ranges in large_chunks]
        metadata_tracker.update_file_metadata(fpath, chunking_mode, total_chunks, chunk_sizes)
        
        files_indexed += 1
        print()
    
    # Note: ChromaDB 1.0+ automatically persists data when using persist_directory
    print_step("Data persistence...")
    print_success("ChromaDB data will be automatically persisted")
    
    print_header("Indexing Complete!")
    print_success(f"Total files indexed: {files_indexed}")
    if ENABLE_GRANULAR_CHUNKING:
        print_success(f"Granular chunks: {granular_chunks}")
    print_success(f"File-level chunks: {filelevel_chunks}")
    if skipped_files > 0:
        print_warning(f"Skipped files: {skipped_files}")
    print()

def index_folder(folder, chunking_mode='dual'):
    """Legacy wrapper for backward compatibility"""
    return index_single_folder(folder, chunking_mode)

def search_files(query, top_files=None, top_chunks_per_file=None):
    """
    Search function that adapts based on chunking configuration.
    """
    if top_files is None:
        top_files = TOP_FILES_COUNT
    if top_chunks_per_file is None:
        top_chunks_per_file = CHUNKS_PER_FILE
    
    print_header(f"Search Results: '{query}'")
    print_step("Encoding search query...")
    q_emb = model.encode(query)
    print_success("query encoded")
    
    # Step 1: Find top files using file-level chunks
    print_step("Searching file-level chunks for relevant files...")
    try:
        file_results = filelevel_collection.query(query_embeddings=[q_emb], n_results=min(top_files * 2, MAX_SEARCH_RESULTS))  # Get more to account for duplicates
    except Exception as e:
        print_error(f"Error during file-level query: {e}")
        return
    
    if SHOW_DEBUG_OUTPUT:
        print("DEBUG: Raw results dict:")
        print(file_results)
    
    if not file_results["metadatas"][0]:
        print_warning("No results found.")
        return
    
    # Extract unique files from file-level results with best scores
    file_scores = {}
    file_metadata = {}
    
    for i, meta in enumerate(file_results["metadatas"][0]):
        file_path = meta["path"]
        score = file_results["distances"][0][i]
        confidence = 1 - score
        
        if file_path not in file_scores or confidence > file_scores[file_path]:
            file_scores[file_path] = confidence
            file_metadata[file_path] = meta
    
    # Sort files by score and take top N
    sorted_files = sorted(file_scores.items(), key=lambda x: x[1], reverse=True)[:top_files]
    
    print_success(f"Found {len(sorted_files)} relevant files")
    print()
    
    # Step 2: For each file, show results based on chunking mode
    for file_idx, (file_path, file_confidence) in enumerate(sorted_files):
        meta = file_metadata[file_path]
        fname = meta["fname"]
        file_type = meta["file_type"]
        
        print(f"{Colors.BOLD}{file_idx+1}. {fname} [{file_type}]{Colors.ENDC} {Colors.OKGREEN}(File Score: {file_confidence:.2f}){Colors.ENDC}")
        print(f"   {Colors.OKCYAN}Path: {file_path}{Colors.ENDC}")
        
        if ENABLE_GRANULAR_CHUNKING:
            # Dual-chunk mode: Show top lines within this file
            try:
                line_results = granular_collection.query(
                    query_embeddings=[q_emb], 
                    n_results=top_chunks_per_file,
                    where={"path": file_path}
                )
                
                if line_results["metadatas"][0]:
                    for j, line_meta in enumerate(line_results["metadatas"][0]):
                        line_num = line_meta["line_num"]
                        line_doc = line_results["documents"][0][j]
                        print(f"      {Colors.WARNING}Line {line_num}:{Colors.ENDC} {line_doc.strip()}")
                else:
                    print(f"      {Colors.WARNING}No specific lines found in this file{Colors.ENDC}")
                    
            except Exception as e:
                print_error(f"Error searching lines in file: {e}")
        else:
            # File-level only mode: Show top chunks within this file
            try:
                chunk_results = filelevel_collection.query(
                    query_embeddings=[q_emb], 
                    n_results=top_chunks_per_file,
                    where={"path": file_path}
                )
                
                if chunk_results["metadatas"][0]:
                    for j, chunk_meta in enumerate(chunk_results["metadatas"][0]):
                        chunk_id = chunk_meta["chunk_id"]
                        chunk_doc = chunk_results["documents"][0][j]
                        chunk_size = chunk_meta["chunk_size"]
                        # Truncate long chunks for display
                        display_chunk = chunk_doc[:MAX_CHUNK_DISPLAY_LENGTH] + "..." if len(chunk_doc) > MAX_CHUNK_DISPLAY_LENGTH else chunk_doc
                        print(f"      {Colors.WARNING}Chunk {chunk_id} ({chunk_size} lines):{Colors.ENDC}")
                        print(f"        {display_chunk.strip()}")
                        print()
                else:
                    print(f"      {Colors.WARNING}No chunks found in this file{Colors.ENDC}")
                    
            except Exception as e:
                print_error(f"Error searching chunks in file: {e}")
        
        print()

def select_folders_to_index():
    """Interactive folder selection for CLI"""
    print_header("folder selection")
    
    # Get currently saved folders
    saved_folders = folder_settings.get_indexed_folders()
    
    if saved_folders:
        print_step("currently indexed folders:")
        for i, folder in enumerate(saved_folders, 1):
            print(f"  {i}. {folder}")
        print()
    
    while True:
        if saved_folders:
            print(f"{Colors.BOLD}options:{Colors.ENDC}")
            print("1. use current folders")
            print("2. add new folder")  
            print("3. remove folder")
            print("4. clear all folders")
            print("5. done")
        else:
            print("no folders selected yet. please add folders to index.")
            print("1. add folder")
            print("2. exit")
        
        try:
            choice = input("\nchoice: ").strip()
            
            if not saved_folders:
                if choice == "1":
                    # Add folder
                    folder_path = input("enter folder path: ").strip()
                    if folder_path:
                        folder_path = os.path.abspath(os.path.expanduser(folder_path))
                        if os.path.exists(folder_path) and os.path.isdir(folder_path):
                            saved_folders = [folder_path]
                            folder_settings.save_folder_settings(saved_folders, [], 10, 'gist')
                            print_success(f"added: {folder_path}")
                        else:
                            print_error("folder does not exist")
                elif choice == "2":
                    return []
            else:
                if choice == "1":
                    return saved_folders
                elif choice == "2":
                    # Add folder
                    folder_path = input("enter folder path: ").strip()
                    if folder_path:
                        folder_path = os.path.abspath(os.path.expanduser(folder_path))
                        if os.path.exists(folder_path) and os.path.isdir(folder_path):
                            if folder_path not in saved_folders:
                                saved_folders.append(folder_path)
                                folder_settings.save_folder_settings(saved_folders, [], 10, 'gist')
                                print_success(f"added: {folder_path}")
                            else:
                                print_warning("folder already in list")
                        else:
                            print_error("folder does not exist")
                elif choice == "3":
                    # Remove folder
                    if len(saved_folders) == 1:
                        print_warning("cannot remove the only folder")
                        continue
                    print("select folder to remove:")
                    for i, folder in enumerate(saved_folders, 1):
                        print(f"  {i}. {folder}")
                    try:
                        idx = int(input("folder number: ").strip()) - 1
                        if 0 <= idx < len(saved_folders):
                            removed = saved_folders.pop(idx)
                            folder_settings.save_folder_settings(saved_folders, [], 10, 'gist')
                            print_success(f"removed: {removed}")
                        else:
                            print_error("invalid selection")
                    except ValueError:
                        print_error("please enter a number")
                elif choice == "4":
                    # Clear all
                    confirm = input("clear all folders? (y/N): ").strip().lower()
                    if confirm == 'y':
                        saved_folders = []
                        folder_settings.save_folder_settings([], [], 10, 'gist')
                        print_success("cleared all folders")
                elif choice == "5":
                    return saved_folders
                else:
                    print_error("invalid choice")
        except KeyboardInterrupt:
            return []

def select_chunking_mode():
    """Select chunking mode for indexing"""
    print_header("chunking mode selection")
    print("1. gist mode (recommended) - better for topic/document search")
    print("2. pinpoint mode - better for exact phrase matching")
    
    while True:
        try:
            choice = input("\nselect mode (1 or 2): ").strip()
            if choice == "1":
                return "gist"
            elif choice == "2":
                return "pinpoint"
            else:
                print_error("please enter 1 or 2")
        except KeyboardInterrupt:
            return "gist"

def show_main_menu():
    """Show the main CLI menu"""
    print(f"\n{Colors.BOLD}filefinder cli{Colors.ENDC}")
    print("1. index folders")
    print("2. search indexed content")
    print("3. manage folders")
    print("4. change chunking mode")
    print("5. show indexed folders")
    print("6. exit")

if __name__ == "__main__":
    # Show supported formats
    show_supported_formats()
    
    # Get folders and chunking mode
    folders = folder_settings.get_indexed_folders()
    current_mode = "gist"  # default mode
    
    # Main CLI loop
    while True:
        try:
            show_main_menu()
            choice = input("\nchoice: ").strip()
            
            if choice == "1":
                # Index folders
                if not folders:
                    print_warning("no folders selected. please select folders first.")
                    folders = select_folders_to_index()
                    if not folders:
                        continue
                
                # Select chunking mode
                current_mode = select_chunking_mode()
                
                # Index the folders
                index_folders(folders, current_mode)
                
            elif choice == "2":
                # Search indexed content
                if not folders:
                    print_warning("no folders indexed yet. please index folders first.")
                    continue
                
                print_header("search mode")
                print("enter search terms (type 'back' to return to menu)")
                
                while True:
                    query = input("\nsearch: ").strip()
                    if query.lower() in ['back', 'exit', 'menu']:
                        break
                    if query:
                        search_files(query)
                
            elif choice == "3":
                # Manage folders
                folders = select_folders_to_index()
                
            elif choice == "4":
                # Change chunking mode
                current_mode = select_chunking_mode()
                print_success(f"chunking mode set to: {current_mode}")
                
            elif choice == "5":
                # Show indexed folders
                if folders:
                    print_header("indexed folders")
                    for i, folder in enumerate(folders, 1):
                        print(f"  {i}. {folder}")
                    print(f"\ncurrent chunking mode: {current_mode}")
                else:
                    print_warning("no folders indexed")
                
            elif choice == "6":
                # Exit
                print("goodbye!")
                break
                
            else:
                print_error("invalid choice, please try again")
                
        except KeyboardInterrupt:
            print("\ngoodbye!")
            break
        except Exception as e:
            print_error(f"unexpected error: {e}")
            continue
